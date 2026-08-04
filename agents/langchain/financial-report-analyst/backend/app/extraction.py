from __future__ import annotations

import csv
import io
from pathlib import Path
from typing import Any

import pandas as pd

from app.models import json_dumps


def rows_to_markdown(rows: list[list[Any]]) -> str:
    cleaned = [["" if value is None else str(value).strip() for value in row] for row in rows if any(row)]
    if not cleaned:
        return ""
    width = max(len(row) for row in cleaned)
    normalized = [row + [""] * (width - len(row)) for row in cleaned]
    header = normalized[0]
    lines = ["| " + " | ".join(header) + " |", "| " + " | ".join(["---"] * width) + " |"]
    for row in normalized[1:]:
        lines.append("| " + " | ".join(row) + " |")
    return "\n".join(lines)


def extract_pdf(path: Path, artifact_dir: Path) -> list[dict[str, Any]]:
    import fitz
    import pdfplumber

    pages: list[dict[str, Any]] = []
    doc = fitz.open(path)
    tables_by_page: dict[int, list[list[list[Any]]]] = {}
    try:
        with pdfplumber.open(path) as pdf:
            for index, page in enumerate(pdf.pages, start=1):
                tables_by_page[index] = page.extract_tables() or []
    except Exception:
        tables_by_page = {}

    for index, page in enumerate(doc, start=1):
        text = page.get_text("text") or ""
        image_path = artifact_dir / "pages" / f"page-{index}.png"
        pix = page.get_pixmap(matrix=fitz.Matrix(1.6, 1.6), alpha=False)
        image_path.parent.mkdir(parents=True, exist_ok=True)
        pix.save(image_path)
        pages.append(
            {
                "page_number": index,
                "text": text,
                "image_path": str(image_path),
                "tables": [
                    {"rows": table, "rows_json": json_dumps(table), "table_markdown": rows_to_markdown(table)}
                    for table in tables_by_page.get(index, [])
                    if rows_to_markdown(table)
                ],
            }
        )
    return pages


def extract_pptx(path: Path, artifact_dir: Path) -> list[dict[str, Any]]:
    from pptx import Presentation

    prs = Presentation(path)
    pages: list[dict[str, Any]] = []
    for index, slide in enumerate(prs.slides, start=1):
        chunks: list[str] = []
        tables: list[dict[str, str]] = []
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False):
                chunks.append(shape.text)
            if getattr(shape, "has_table", False):
                rows = [[cell.text for cell in row.cells] for row in shape.table.rows]
                markdown = rows_to_markdown(rows)
                if markdown:
                    tables.append({"rows": rows, "rows_json": json_dumps(rows), "table_markdown": markdown})
        pages.append({"page_number": index, "text": "\n".join(chunks), "image_path": "", "tables": tables})
    return pages


def extract_docx(path: Path, artifact_dir: Path) -> list[dict[str, Any]]:
    from docx import Document

    doc = Document(path)
    text = "\n".join(paragraph.text for paragraph in doc.paragraphs)
    tables: list[dict[str, str]] = []
    for table in doc.tables:
        rows = [[cell.text for cell in row.cells] for row in table.rows]
        markdown = rows_to_markdown(rows)
        if markdown:
            tables.append({"rows": rows, "rows_json": json_dumps(rows), "table_markdown": markdown})
    return [{"page_number": 1, "text": text, "image_path": "", "tables": tables}]


def extract_spreadsheet(path: Path, artifact_dir: Path) -> list[dict[str, Any]]:
    sheets = pd.read_excel(path, sheet_name=None) if path.suffix.lower() in {".xlsx", ".xls"} else {"CSV": pd.read_csv(path)}
    tables: list[dict[str, str]] = []
    text_chunks: list[str] = []
    for name, frame in sheets.items():
        frame = frame.fillna("")
        rows = [list(frame.columns)] + frame.astype(str).values.tolist()
        markdown = rows_to_markdown(rows[:60])
        if markdown:
            tables.append({"rows": rows, "rows_json": frame.to_json(orient="records"), "table_markdown": f"### {name}\n{markdown}"})
            text_chunks.append(f"Sheet {name}\n{frame.to_csv(index=False)}")
    return [{"page_number": 1, "text": "\n".join(text_chunks), "image_path": "", "tables": tables}]


def extract_text_file(path: Path, artifact_dir: Path) -> list[dict[str, Any]]:
    if path.suffix.lower() == ".csv":
        with path.open(newline="") as csv_file:
            rows = list(csv.reader(csv_file))
        return [{"page_number": 1, "text": path.read_text(errors="ignore"), "image_path": "", "tables": [{"rows": rows, "rows_json": json_dumps(rows), "table_markdown": rows_to_markdown(rows)}]}]
    return [{"page_number": 1, "text": path.read_text(errors="ignore"), "image_path": "", "tables": []}]


def extract_image(path: Path, artifact_dir: Path) -> list[dict[str, Any]]:
    target = artifact_dir / "pages" / "page-1.png"
    target.parent.mkdir(parents=True, exist_ok=True)
    from PIL import Image

    image = Image.open(path)
    image.thumbnail((1800, 1800))
    image.convert("RGB").save(target)
    return [{"page_number": 1, "text": "", "image_path": str(target), "tables": []}]


def extract_document(path: Path, artifact_dir: Path) -> list[dict[str, Any]]:
    suffix = path.suffix.lower()
    if suffix == ".pdf":
        return extract_pdf(path, artifact_dir)
    if suffix == ".pptx":
        return extract_pptx(path, artifact_dir)
    if suffix == ".docx":
        return extract_docx(path, artifact_dir)
    if suffix in {".xlsx", ".xls", ".csv"}:
        return extract_spreadsheet(path, artifact_dir)
    if suffix in {".png", ".jpg", ".jpeg", ".webp"}:
        return extract_image(path, artifact_dir)
    return extract_text_file(path, artifact_dir)
